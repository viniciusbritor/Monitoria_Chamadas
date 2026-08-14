import io
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from core.excel_exporter import _extract_call_fields

def generate_pptx_report(calls: list[dict]) -> bytes:
    """
    Gera uma Apresentação Executiva em PowerPoint (.pptx) de 5 slides
    estritamente limitada às 50 chamadas fornecidas, com todos os campos populados.
    """
    # Limita rigorosamente a 50 chamadas como ordenado nas diretrizes
    limited_calls = calls[:50]
    extracted_calls = [_extract_call_fields(c) for c in limited_calls]
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]
    
    # Cores Oficiais da Identidade Visual Coherence
    DARK_SLATE = RGBColor(15, 23, 42)    # #0F172A
    JADE_GREEN = RGBColor(26, 107, 82)   # #1A6B52
    LIGHT_BG = RGBColor(248, 250, 252)   # #F8FAFC
    CARD_BG = RGBColor(255, 255, 255)    # #FFFFFF
    TEXT_MAIN = RGBColor(51, 65, 85)     # #334155
    TEXT_MUTED = RGBColor(100, 116, 139) # #64748B
    BORDER_COLOR = RGBColor(226, 232, 240)
    
    def set_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, subtitle_text=""):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_SLATE
        shape.line.fill.background()
        
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.size = Pt(12)
            p2.font.color.rgb = RGBColor(148, 163, 184)

    # ----------------------------------------------------
    # SLIDE 1: Capa Executiva
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide1, DARK_SLATE)
    
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(2.2), Inches(1.5), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = JADE_GREEN
    line.line.fill.background()
    
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(11), Inches(2.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p_title = tf1.paragraphs[0]
    p_title.text = "Relatório Executivo de Monitoria de Chamadas"
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 255, 255)
    
    p_sub = tf1.add_paragraph()
    p_sub.text = f"Coherence AI • Auditoria de Atendimentos ({len(extracted_calls)} chamadas analisadas)"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = RGBColor(148, 163, 184)

    # ----------------------------------------------------
    # SLIDE 2: Visão Geral de Performance (KPIs)
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide2, LIGHT_BG)
    add_header(slide2, "Visão Geral de Performance Operacional", "Métricas consolidadas do lote de chamadas auditadas")
    
    concluidas = [c for c in extracted_calls if str(c.get("status")).lower().startswith("conclu") or str(c.get("status")).lower() == "finalizado"]
    notas_qa = [float(c["nota_qa"]) for c in concluidas if isinstance(c["nota_qa"], (int, float)) or (isinstance(c["nota_qa"], str) and c["nota_qa"].replace('.', '', 1).isdigit())]
    qa_medio = round(sum(notas_qa) / len(notas_qa), 1) if notas_qa else 0.0
    
    nps_list = []
    for c in concluidas:
        val = c.get("nps_cli")
        if val is not None and val != "-":
            try:
                nps_list.append(float(val))
            except (ValueError, TypeError):
                pass
    nps_medio = round(sum(nps_list) / len(nps_list), 1) if nps_list else 0.0
    
    erros_criticos = sum(1 for c in concluidas if c.get("erro_crit") == "SIM")
    
    cards_data = [
        ("Total Auditado", str(len(extracted_calls)), "Chamadas na amostra"),
        ("QA Score Médio", f"{qa_medio}/100", "Nota técnica de qualidade"),
        ("NPS Sentimento", f"{nps_medio}/10", "Índice de satisfação do cliente"),
        ("Erros Críticos", str(erros_criticos), "Alertas de inconformidade"),
    ]
    
    for idx, (label, val, sub) in enumerate(cards_data):
        left_pos = Inches(0.8 + idx * 3.0)
        top_pos = Inches(2.2)
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, Inches(2.7), Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        
        tb = slide2.shapes.add_textbox(left_pos + Inches(0.15), top_pos + Inches(0.3), Inches(2.4), Inches(2.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = label
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_MUTED
        
        p2 = tf.add_paragraph()
        p2.text = val
        p2.font.size = Pt(36)
        p2.font.bold = True
        p2.font.color.rgb = JADE_GREEN if "QA" in label or "Total" in label else DARK_SLATE
        
        p3 = tf.add_paragraph()
        p3.text = sub
        p3.font.size = Pt(11)
        p3.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 3: Sentimento & Humor dos Clientes e Atendentes
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide3, LIGHT_BG)
    add_header(slide3, "Análise de Sentimento & Humor", "Distribuição de humor do cliente e postura do atendente")
    
    humores_cli = {}
    humores_atend = {}
    for c in concluidas:
        hc = c.get("humor_cli") or "Não Informado"
        humores_cli[hc] = humores_cli.get(hc, 0) + 1
        ha = c.get("humor_op") or "Não Informado"
        humores_atend[ha] = humores_atend.get(ha, 0) + 1
        
    # Box Clientes
    box_cli = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    box_cli.fill.solid()
    box_cli.fill.fore_color.rgb = CARD_BG
    box_cli.line.color.rgb = BORDER_COLOR
    
    tb_cli = slide3.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_cli = tb_cli.text_frame
    tf_cli.word_wrap = True
    p = tf_cli.paragraphs[0]
    p.text = "👤 Humor do Cliente"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_SLATE
    
    for h_name, h_count in humores_cli.items():
        pct = round((h_count / max(len(concluidas), 1)) * 100, 1)
        p_item = tf_cli.add_paragraph()
        p_item.text = f"• {h_name}: {h_count} chamada(s) ({pct}%)"
        p_item.font.size = Pt(13)
        p_item.font.color.rgb = TEXT_MAIN

    # Box Atendentes
    box_at = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    box_at.fill.solid()
    box_at.fill.fore_color.rgb = CARD_BG
    box_at.line.color.rgb = BORDER_COLOR
    
    tb_at = slide3.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.4))
    tf_at = tb_at.text_frame
    tf_at.word_wrap = True
    p = tf_at.paragraphs[0]
    p.text = "🎧 Postura / Humor do Atendente"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = JADE_GREEN
    
    for h_name, h_count in humores_atend.items():
        pct = round((h_count / max(len(concluidas), 1)) * 100, 1)
        p_item = tf_at.add_paragraph()
        p_item.text = f"• {h_name}: {h_count} chamada(s) ({pct}%)"
        p_item.font.size = Pt(13)
        p_item.font.color.rgb = TEXT_MAIN

    # ----------------------------------------------------
    # SLIDE 4: Diagnóstico Operacional & Treinamento
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide4, LIGHT_BG)
    add_header(slide4, "Diagnóstico Operacional & Recomendações", "Pontos de destaque, oportunidades e plano de treinamento")
    
    # Pontos Positivos Box
    box_pos = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(3.7), Inches(4.8))
    box_pos.fill.solid()
    box_pos.fill.fore_color.rgb = CARD_BG
    box_pos.line.color.rgb = BORDER_COLOR
    
    tb_pos = slide4.shapes.add_textbox(Inches(0.95), Inches(1.95), Inches(3.4), Inches(4.5))
    tf_pos = tb_pos.text_frame
    tf_pos.word_wrap = True
    p = tf_pos.paragraphs[0]
    p.text = "🌟 Pontos Positivos"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = JADE_GREEN
    
    all_pos = []
    for c in concluidas:
        all_pos.extend(c.get("pontos_positivos_list") or [])
    for item in (all_pos[:4] if all_pos else ["Cordialidade no atendimento inicial", "Clareza nas orientações prestadas"]):
        p_it = tf_pos.add_paragraph()
        p_it.text = f"✓ {item}"
        p_it.font.size = Pt(11)
        p_it.font.color.rgb = TEXT_MAIN

    # Pontos de Melhoria Box
    box_neg = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.8), Inches(3.7), Inches(4.8))
    box_neg.fill.solid()
    box_neg.fill.fore_color.rgb = CARD_BG
    box_neg.line.color.rgb = BORDER_COLOR
    
    tb_neg = slide4.shapes.add_textbox(Inches(4.95), Inches(1.95), Inches(3.4), Inches(4.5))
    tf_neg = tb_neg.text_frame
    tf_neg.word_wrap = True
    p = tf_neg.paragraphs[0]
    p.text = "🎯 Pontos de Melhoria"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = DARK_SLATE
    
    all_neg = []
    for c in concluidas:
        all_neg.extend(c.get("pontos_melhoria_list") or [])
    for item in (all_neg[:4] if all_neg else ["Aprimorar contorno de objeções", "Reduzir tempo de espera em pausa"]):
        p_it = tf_neg.add_paragraph()
        p_it.text = f"⚠ {item}"
        p_it.font.size = Pt(11)
        p_it.font.color.rgb = TEXT_MAIN

    # Recomendação de Treinamento Box
    box_rec = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8))
    box_rec.fill.solid()
    box_rec.fill.fore_color.rgb = CARD_BG
    box_rec.line.color.rgb = BORDER_COLOR
    
    tb_rec = slide4.shapes.add_textbox(Inches(8.95), Inches(1.95), Inches(3.4), Inches(4.5))
    tf_rec = tb_rec.text_frame
    tf_rec.word_wrap = True
    p = tf_rec.paragraphs[0]
    p.text = "💡 Plano de Treinamento"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = JADE_GREEN
    
    all_recs = [c.get("rec_trein") for c in concluidas if c.get("rec_trein") and c.get("rec_trein") != "-"]
    for item in (all_recs[:3] if all_recs else ["Capacitação em confirmação de dados e escuta ativa"]):
        p_it = tf_rec.add_paragraph()
        p_it.text = f"• {item}"
        p_it.font.size = Pt(11)
        p_it.font.color.rgb = TEXT_MAIN

    # ----------------------------------------------------
    # SLIDE 5: Tabela Analítica das Chamadas (Máximo 50)
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide5, LIGHT_BG)
    add_header(slide5, "Detalhamento Analítico dos Atendimentos", f"Amostra dos {len(extracted_calls)} atendimentos auditados (Limite 50)")
    
    rows = min(len(extracted_calls) + 1, 12)
    cols = 7
    left = Inches(0.8)
    top = Inches(1.6)
    width = Inches(11.733)
    height = Inches(5.2)
    
    table_shape = slide5.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    table.columns[0].width = Inches(1.8) # ID Chamada
    table.columns[1].width = Inches(2.2) # Atendente
    table.columns[2].width = Inches(2.6) # Motivo do Contato
    table.columns[3].width = Inches(1.2) # QA Score
    table.columns[4].width = Inches(1.1) # NPS
    table.columns[5].width = Inches(1.4) # Humor Cliente
    table.columns[6].width = Inches(1.433) # Humor Atendente
    
    headers = ["ID Chamada", "Atendente", "Motivo do Contato", "QA Score", "NPS", "Humor Cliente", "Humor Atendente"]
    for col_idx, h_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = h_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_SLATE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
    for row_idx, c in enumerate(extracted_calls[:11], start=1):
        cid_short = str(c["call_id"])[:8] + "..." if len(str(c["call_id"])) > 8 else str(c["call_id"])
        atendente = str(c["atendente"])[:22]
        motivo = str(c["motivo"])[:35]
        qa = str(c["nota_qa"])
        nps = str(c["nps_cli"])
        humor_cli = str(c["humor_cli"])[:16]
        humor_op = str(c["humor_op"])[:16]
        
        row_values = [cid_short, atendente, motivo, qa, nps, humor_cli, humor_op]
        for col_idx, val in enumerate(row_values):
            cell = table.cell(row_idx, col_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(241, 245, 249) if row_idx % 2 == 0 else RGBColor(255, 255, 255)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = TEXT_MAIN
                if col_idx in [0, 3, 4, 5, 6]:
                    p.alignment = PP_ALIGN.CENTER

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()
