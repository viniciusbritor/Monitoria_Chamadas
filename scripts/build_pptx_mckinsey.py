"""
Build McKinsey-style PPTX using shared cost model.
All values from core/cost_model.py — guaranteed consistent with Excel.
"""
import json, os, sys
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from core.cost_model import *

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUT = r"C:\Users\vinic\workspace_antigravity\Monitoria_Chamadas\docs\Custos_Projecao_OmniChannel.pptx"
LOGO = r"C:\Users\vinic\workspace_antigravity\Monitoria_Chamadas\frontend\public\logo-top-v2.png"
QR = r"C:\Users\vinic\workspace_antigravity\Monitoria_Chamadas\docs\pix_qr.png"
FID = "1Nb_OLbJS0012keYcXW58EMz4F1evMz6w"  # Custos/ dentro de Omnichannel/
from datetime import date
DT = date.today().isoformat()
FNAME = f"Custos_Projecao_OmniChannel_{DT}.pptx"
TKN = os.path.expanduser(r"~\.gemini\config\skills\google_calendar_manager\resources\token_drive.json")

# McKinsey colors
NAVY = RGBColor(0x00, 0x2B, 0x4D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GREEN_C = RGBColor(0x00, 0x61, 0x00)
BLUE_C = RGBColor(0x00, 0x5B, 0x96)
ORANGE_C = RGBColor(0xD3, 0x54, 0x00)
MID_BLUE = RGBColor(0x00, 0x73, 0xBE)
LIGHT_BLUE_BG = RGBColor(0xD6, 0xE8, 0xF7)
TB = RGBColor(0xDD, 0xDD, 0xDD)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H

def hex_color(h):
    return RGBColor(int(h[1:3],16), int(h[3:5],16), int(h[5:7],16))

def add_blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def tb(slide, l, t, w, h, text, size=12, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = 'Calibri'; p.alignment = align
    return box

def ap(tf, text, size=11, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT, sb=4):
    p_ = tf.add_paragraph(); p_.text = text; p_.font.size = Pt(size); p_.font.bold = bold
    p_.font.color.rgb = color; p_.font.name = 'Calibri'; p_.alignment = align; p_.space_before = Pt(sb)
    return p_

def at(slide, l, t, w, h, rows, cols):
    return slide.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w), Inches(h)).table

def sc(tbl, r, c, text, size=10, bold=False, color=None, bg=None, align=PP_ALIGN.CENTER):
    cell = tbl.cell(r, c); cell.text = ""
    p_ = cell.text_frame.paragraphs[0]; p_.text = str(text); p_.font.size = Pt(size)
    p_.font.bold = bold; p_.font.name = 'Calibri'; p_.alignment = align
    if color: p_.font.color.rgb = color
    if bg: cell.fill.solid(); cell.fill.fore_color.rgb = bg
    cell.margin_left = Pt(4); cell.margin_right = Pt(4); cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
    return cell

def logo(slide):
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(SLIDE_W.inches - 2.0), Inches(SLIDE_H.inches - 0.65), height=Inches(0.4))

def footer(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(SLIDE_H.inches - 0.55), SLIDE_W, Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = LIGHT_GRAY; bar.line.fill.background()
    tb(slide, 0.5, SLIDE_H.inches - 0.45, 4, 0.3, "CONFIDENCIAL - OmniChannel 2026", 7, color=MED_GRAY)
    tb(slide, SLIDE_W.inches - 3.5, SLIDE_H.inches - 0.45, 3, 0.3, "Fonte: GCP Pricing Jul/2026, DeepSeek API", 7, color=MED_GRAY, align=PP_ALIGN.RIGHT)
    logo(slide)

def divider(slide, top):
    d = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top), Inches(11.7), Pt(1))
    d.fill.solid(); d.fill.fore_color.rgb = TB; d.line.fill.background()

def pn(slide, n):
    tb(slide, SLIDE_W.inches/2 - 0.5, SLIDE_H.inches - 0.5, 1, 0.4, str(n), 9, color=MED_GRAY, align=PP_ALIGN.CENTER)

# ═══════ SLIDE 1: CAPA ═══════
s = add_blank()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.8), SLIDE_W, Inches(0.05))
bar.fill.solid(); bar.fill.fore_color.rgb = MID_BLUE; bar.line.fill.background()
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(1.2), Inches(1.2), height=Inches(0.8))
tb(s, 1.2, 3.2, 10, 1.2, "CUSTOS OPERACIONAIS", 40, True, WHITE)
tb(s, 1.2, 4.1, 10, 0.8, "Infraestrutura & Inteligencia Artificial", 28, color=RGBColor(0xAA, 0xCC, 0xEE))
tb(s, 1.2, 5.2, 10, 0.5, "Plataforma OmniChannel - Julho 2026", 14, color=MED_GRAY)
tb(s, 1.2, 6.5, 10, 0.4, "CONFIDENCIAL - Preparado para investidores e stakeholders", 9, color=MED_GRAY)

# ═══════ SLIDE 2: PREMISSAS ═══════
s = add_blank()
tb(s, 0.8, 0.5, 11, 0.7, "OPERACAO EM GCP COM DEEPSEEK V4 FLASH COMO LLM PRIMARIO", 22, True, NAVY)
divider(s, 1.2); footer(s); pn(s, 2)
tb(s, 0.8, 1.5, 11, 0.4, "Premissas financeiras e tecnicas que fundamentam todas as projecoes", 12, color=MED_GRAY)

t1 = at(s, 0.8, 2.2, 11.5, 4.0, len(PREMISSAS)+1, 3)
for r, (a, b, c) in enumerate([("Premissa", "Valor", "Fonte")] + PREMISSAS):
    hdr = r == 0; bg_c = NAVY if hdr else (LIGHT_GRAY if r%2==0 else None)
    sc(t1, r, 0, a, 11, hdr, WHITE if hdr else DARK_GRAY, bg_c, PP_ALIGN.LEFT)
    sc(t1, r, 1, b, 11, hdr, WHITE if hdr else DARK_GRAY, bg_c)
    sc(t1, r, 2, c, 10, hdr, WHITE if hdr else MED_GRAY, bg_c, PP_ALIGN.LEFT)
t1.columns[0].width = Inches(3.5); t1.columns[1].width = Inches(4.0); t1.columns[2].width = Inches(4.0)

# ═══════ SLIDE 3: 3 PILARES ═══════
s = add_blank()
tb(s, 0.8, 0.5, 11, 0.7, f"R$ {TOTAL_HOJE*BRL:,.0f}/MES: GCP {GCP_TOTAL/TOTAL_HOJE*100:.0f}%, IA DEV {AI_DEV_TOTAL/TOTAL_HOJE*100:.0f}%, IA PROD {AI_PROD_TOTAL/TOTAL_HOJE*100:.0f}%", 22, True, NAVY)
divider(s, 1.2); footer(s); pn(s, 3)
tb(s, 0.8, 1.5, 11, 0.4, "Tres pilares independentes de custo operacional", 12, color=MED_GRAY)

cards = [
    ("INFRAESTRUTURA GCP", GCP_TOTAL, "Cloud Run, Firestore, Storage,\nPub/Sub, Compute Engine", NAVY),
    ("IA DESENVOLVIMENTO", AI_DEV_TOTAL, "DeepSeek tokens para testes,\nengenharia de prompts,\nMiniMax Plus (plano fixo)", MID_BLUE),
    ("IA PRODUCAO (MODULOS)", AI_PROD_TOTAL, "DeepSeek na Monitoria,\nChatbots WhatsApp,\nextras (URA, overflow)", GREEN_C),
]
for i, (title, usd, desc, clr) in enumerate(cards):
    x = 0.8 + i * 4.0
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(3.8), Inches(4.5))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT_GRAY; card.line.fill.background()
    acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.2), Inches(3.8), Pt(6))
    acc.fill.solid(); acc.fill.fore_color.rgb = clr; acc.line.fill.background()
    tb(s, x+0.3, 2.5, 3.2, 0.5, title, 14, True, clr, PP_ALIGN.CENTER)
    tb(s, x+0.3, 3.2, 3.2, 0.8, f"R$ {usd*BRL:,.0f}", 32, True, DARK_GRAY, PP_ALIGN.CENTER)
    tb(s, x+0.3, 3.9, 3.2, 0.4, f"USD ${usd:,.0f}/mes", 12, color=MED_GRAY, align=PP_ALIGN.CENTER)
    tb(s, x+0.3, 4.6, 3.2, 1.8, desc, 10, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ═══════ SLIDE 4: CUSTOS HOJE ═══════
s = add_blank()
tb(s, 0.8, 0.5, 11, 0.7,
   f"GCP R$ {GCP_TOTAL*BRL:,.0f} + IA DESENVOLVIMENTO R$ {AI_DEV_TOTAL*BRL:,.0f} + IA PRODUCAO R$ {AI_PROD_TOTAL*BRL:,.0f} = R$ {TOTAL_HOJE*BRL:,.0f}/MES",
   22, True, NAVY)
divider(s, 1.2); footer(s); pn(s, 4)

for si, (title, data, total, color, x_pos) in enumerate([
    ("Infraestrutura GCP", INFRA, GCP_TOTAL, NAVY, 0.8),
    ("IA - Desenvolvimento", AI_DEV, AI_DEV_TOTAL, MID_BLUE, 7.0),
    ("IA - Producao (Modulos)", AI_PROD, AI_PROD_TOTAL, GREEN_C, 10.3),
]):
    tb(s, x_pos, 1.5, 5, 0.4, title, 14, True, color)
    w = 5.5 if si == 0 else (5.0 if si == 1 else 2.5)
    t = at(s, x_pos, 2.0, w, 0.35*(len(data)+1)+0.4, len(data)+2, 2)
    sc(t, 0, 0, "Servico" if si==0 else "Uso de IA", 9, True, WHITE, color, PP_ALIGN.LEFT)
    sc(t, 0, 1, "BRL/mes", 9, True, WHITE, color)
    for r, (nm, usd) in enumerate(data, 1):
        sc(t, r, 0, nm, 9, align=PP_ALIGN.LEFT)
        sc(t, r, 1, f"R$ {usd*BRL:,.0f}", 9)
    sc(t, len(data)+1, 0, "Subtotal", 10, True, DARK_GRAY, LIGHT_BLUE_BG, PP_ALIGN.LEFT)
    sc(t, len(data)+1, 1, f"R$ {total*BRL:,.0f}", 10, True, DARK_GRAY, LIGHT_BLUE_BG)
    t.columns[0].width = Inches(4.0 if si==0 else (3.8 if si==1 else 1.5))
    t.columns[1].width = Inches(1.5 if si==0 else (1.2 if si==1 else 1.0))

tb(s, 0.8, 5.5, 11.5, 0.5, f"TOTAL MENSAL: R$ {TOTAL_HOJE*BRL:,.0f}", 24, True, NAVY, PP_ALIGN.CENTER)

# ═══════ SLIDE 5: 500 CHAMADAS/DIA ═══════
s = add_blank()
sc5 = SCENARIOS[500]
tb(s, 0.8, 0.5, 11, 0.7, f"CENARIO-ALVO: 500 CHAMADAS/DIA CUSTA APENAS R$ {sc5['cost_per_call_brl']:.2f} POR CHAMADA", 22, True, NAVY)
divider(s, 1.2); footer(s); pn(s, 5)
tb(s, 0.8, 1.5, 11, 0.4, f"Primeiro cliente BPO - {sc5['calls_month']:,} chamadas/mes processadas automaticamente", 12, color=MED_GRAY)

# Big number card
card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(5.5), Inches(1.5))
card.fill.solid(); card.fill.fore_color.rgb = NAVY; card.line.fill.background()
tb(s, 1.2, 2.4, 4.8, 0.5, f"R$ {sc5['cost_per_call_brl']:.2f}", 48, True, WHITE, PP_ALIGN.CENTER)
tb(s, 1.2, 3.2, 4.8, 0.3, "por chamada analisada", 14, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

# Breakdown table
t3 = at(s, 7.0, 2.2, 6.0, 4.0, 10, 3)
rows = [
    ("Categoria", "BRL/mes", "%"),
    ("GCP (Worker + Infra + WhatsApp)", f"R$ {sc5['gcp']*BRL:,.0f}", f"{sc5['gcp']/sc5['total_usd']*100:.0f}%"),
    ("IA Desenvolvimento (estabilizado)", f"R$ {sc5['ai_dev']*BRL:,.0f}", f"{sc5['ai_dev']/sc5['total_usd']*100:.0f}%"),
    ("IA Producao - DeepSeek Monitoria", f"R$ {sc5['ds_mono']*BRL:,.0f}", f"{sc5['ds_mono']/sc5['total_usd']*100:.0f}%"),
    ("IA Producao - DeepSeek Chatbots", f"R$ {sc5['ds_chat']*BRL:,.0f}", f"{sc5['ds_chat']/sc5['total_usd']*100:.0f}%"),
    ("IA Producao - MiniMax fallback", f"R$ {sc5['mm_token']*BRL:,.0f}", f"{sc5['mm_token']/sc5['total_usd']*100:.0f}%"),
    ("IA Producao - Extras", f"R$ {sc5['ds_extra']*BRL:,.0f}", f"{sc5['ds_extra']/sc5['total_usd']*100:.0f}%"),
    ("", "", ""),
    ("TOTAL 500 chamadas/dia", f"R$ {sc5['total_brl']:,.0f}", "100%"),
    ("CUSTO POR CHAMADA", f"R$ {sc5['cost_per_call_brl']:.2f}", "-"),
]
for r, (a, b, c) in enumerate(rows):
    is_tot = r >= len(rows)-2
    sc(t3, r, 0, a, 12 if is_tot else 9, (r==0 or is_tot), WHITE if r==0 else (GREEN_C if is_tot else DARK_GRAY), NAVY if r==0 else (LIGHT_BLUE_BG if is_tot else None), PP_ALIGN.LEFT)
    sc(t3, r, 1, b, 14 if is_tot else 9, (r==0 or is_tot), WHITE if r==0 else (GREEN_C if is_tot else DARK_GRAY), NAVY if r==0 else (LIGHT_BLUE_BG if is_tot else None))
    sc(t3, r, 2, c, 10, (r==0 or is_tot), WHITE if r==0 else MED_GRAY, NAVY if r==0 else (LIGHT_BLUE_BG if is_tot else None))
t3.columns[0].width = Inches(3.2); t3.columns[1].width = Inches(1.5); t3.columns[2].width = Inches(1.3)

# Scale comparison
tb(s, 0.8, 6.0, 11, 0.3,
   f"Escala: 1.000/dia -> R$ {SCENARIOS[1000]['cost_per_call_brl']:.2f}/chamada  |  5.000/dia -> R$ {SCENARIOS[5000]['cost_per_call_brl']:.2f}/chamada",
   11, color=MED_GRAY, align=PP_ALIGN.CENTER)

# ═══════ SLIDE 6: COMPARATIVO BRASIL ═══════
s = add_blank()
OURS_BRL = sc5['cost_per_call_brl']
tb(s, 0.8, 0.5, 11, 0.7, "TELEPERFORMANCE E ATENTO GASTAM R$ 0,50-1,00/CHAMADA COM QA HUMANO", 22, True, NAVY)
tb(s, 0.8, 1.05, 11, 0.4, f"Nossa automacao com IA e {0.75/OURS_BRL:.0f}-{1.0/OURS_BRL:.0f}x mais barata - e cobre 100% das chamadas", 12, color=MED_GRAY)
divider(s, 1.5); footer(s); pn(s, 6)

t4 = at(s, 1.5, 2.0, 10.3, 3.5, len(BENCHMARKS_BR)+4, 4)
bench_rows = [("Empresa", "Modelo de QA", "Custo por Chamada", "vs Nossa Solucao")]
for emp, perfil, modelo, lo, hi in BENCHMARKS_BR:
    mid = (lo+hi)/2
    bench_rows.append((f"{emp}\n{perfil}", modelo, f"R$ {lo:.2f}-{hi:.2f}", f"{mid/OURS_BRL:.0f}x mais caro"))
bench_rows.append(("NOSSA SOLUCAO", "100% automatizado - IA analisa\nCADA chamada em tempo real\nDeepSeek V4 Flash + MiniMax", f"R$ {OURS_BRL:.2f}", "REFERENCIA"))
bench_rows.append(("Vantagem", f"Cobertura 100% vs 2-5%\nRelatorio em < 2 minutos\nZero intervencao humana", "", ""))

for r, (a, b, c, d) in enumerate(bench_rows):
    is_us = r == len(bench_rows)-2; is_adv = r == len(bench_rows)-1
    bg_c = NAVY if r==0 else (GREEN_C if is_us else (LIGHT_BLUE_BG if is_adv else (LIGHT_GRAY if r%2==0 else None)))
    sc(t4, r, 0, a, 13 if is_us else 10, (r==0 or is_us), WHITE if r==0 or is_us else DARK_GRAY, bg_c, PP_ALIGN.LEFT)
    sc(t4, r, 1, b, 9, color=WHITE if r==0 or is_us else DARK_GRAY, bg=bg_c, align=PP_ALIGN.LEFT)
    sc(t4, r, 2, c, 13 if is_us else 10, (r==0 or is_us), WHITE if r==0 or is_us else DARK_GRAY, bg=bg_c)
    sc(t4, r, 3, d, 10, is_us, WHITE if r==0 or is_us else (GREEN_C if is_us else DARK_GRAY), bg_c)
t4.columns[0].width = Inches(2.8); t4.columns[1].width = Inches(3.0); t4.columns[2].width = Inches(2.0); t4.columns[3].width = Inches(2.5)

savings = int((0.75 - OURS_BRL) * 15000)
tb(s, 0.8, 5.9, 11, 0.4, f"Economia vs Teleperformance: ~R$ {savings:,}/mes para operacao de 500 chamadas/dia", 12, True, GREEN_C, PP_ALIGN.CENTER)

# ═══════ SLIDE 7: CAPACIDADE ═══════
s = add_blank()
tb(s, 0.8, 0.5, 11, 0.7, "PROCESSAMOS 100% DAS CHAMADAS - CONCORRENCIA AUDITA APENAS 2-5%", 22, True, NAVY)
divider(s, 1.2); footer(s); pn(s, 7)
tb(s, 0.8, 1.5, 11, 0.4, "Capacidade atual da plataforma e diferenciais competitivos", 12, color=MED_GRAY)

for i, (num, title, desc) in enumerate(CAPACITY):
    x = 0.8 + i * 4.0
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(3.8), Inches(3.0))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT_GRAY; card.line.fill.background()
    tb(s, x+0.3, 2.4, 3.2, 0.6, num, 36, True, NAVY, PP_ALIGN.CENTER)
    tb(s, x+0.3, 3.1, 3.2, 0.6, title, 12, True, DARK_GRAY, PP_ALIGN.CENTER)
    tb(s, x+0.3, 4.0, 3.2, 1.0, desc, 9, color=MED_GRAY, align=PP_ALIGN.CENTER)

tb(s, 0.8, 5.6, 11, 0.4, "Diferenciais competitivos", 14, True, NAVY)
for i, item in enumerate(DIFFERENTIALS):
    tb(s, 1.2, 6.0 + i*0.22, 11, 0.3, f"- {item}", 9, color=DARK_GRAY)

# ═══════ SLIDE 8: ROADMAP ═══════
s = add_blank()
tb(s, 0.8, 0.5, 11, 0.7, "JULHO 2026 -> ABRIL 2027: DE R$ 1.153/MES A PLATAFORMA SAAS MULTI-CLIENTE", 22, True, NAVY)
divider(s, 1.2); footer(s); pn(s, 8)

for i, (date, phase, action, result, color_str) in enumerate(ROADMAP):
    y = 1.8; x = 0.8 + i * 2.5; clr = hex_color(color_str)
    if i < 4:
        line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+2.3), Inches(y+1.3), Inches(0.4), Pt(3))
        line.fill.solid(); line.fill.fore_color.rgb = MED_GRAY; line.line.fill.background()
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.3), Inches(4.0))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT_GRAY; box.line.fill.background()
    acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(2.3), Pt(6))
    acc.fill.solid(); acc.fill.fore_color.rgb = clr; acc.line.fill.background()
    tb(s, x+0.15, y+0.2, 2.0, 0.3, date, 8, True, clr, PP_ALIGN.CENTER)
    tb(s, x+0.15, y+0.5, 2.0, 0.4, phase, 14, True, DARK_GRAY, PP_ALIGN.CENTER)
    tb(s, x+0.15, y+1.2, 2.0, 1.0, action, 10, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    d2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+0.5), Inches(y+2.5), Inches(1.3), Pt(2))
    d2.fill.solid(); d2.fill.fore_color.rgb = MED_GRAY; d2.line.fill.background()
    tb(s, x+0.15, y+2.6, 2.0, 1.1, result, 9, True, clr, PP_ALIGN.CENTER)

tb(s, 0.8, 6.3, 11.5, 0.4, f"Modelo de receita: cobrar R$ 0,50 por chamada analisada - margem de ~{100-int(OURS_BRL/0.50*100)}% sobre custo operacional de R$ {OURS_BRL:.2f}", 11, True, NAVY, PP_ALIGN.CENTER)

# ═══════ SLIDE 9: PIX QR CODE ═══════
s = add_blank()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(2.0))
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(5.0), Inches(0.3), height=Inches(0.6))
tb(s, 1.0, 2.5, 11, 0.8, "SEGUE MEU PIX", 36, True, NAVY, PP_ALIGN.CENTER)
tb(s, 1.0, 3.3, 11, 0.5, "Apoie o desenvolvimento da plataforma OmniChannel", 16, color=MED_GRAY, align=PP_ALIGN.CENTER)
if os.path.exists(QR):
    s.shapes.add_picture(QR, Inches(4.8), Inches(3.8), height=Inches(2.5))
tb(s, 1.0, 6.5, 11, 0.5, "CPF: 047.799.777-54", 22, True, DARK_GRAY, PP_ALIGN.CENTER)
tb(s, 1.0, 6.95, 11, 0.3, "Vinicius Brito - OmniChannel 2026", 10, color=MED_GRAY, align=PP_ALIGN.CENTER)

# ═══════ SAVE & UPLOAD ═══════
prs.save(OUT)
print(f"PPTX saved: {OUT}")
print(f"  Slides: {len(prs.slides)}")
print(f"  HOJE: R$ {TOTAL_HOJE*BRL:,.0f} (GCP {GCP_TOTAL*BRL:,.0f} + IA Dev {AI_DEV_TOTAL*BRL:,.0f} + IA Prod {AI_PROD_TOTAL*BRL:,.0f})")
for c in [500, 1000, 5000]:
    s_ = SCENARIOS[c]
    print(f"  {c}/dia: R$ {s_['total_brl']:,.0f} | R$ {s_['cost_per_call_brl']:.2f}/chamada")

from google.oauth2.credentials import Credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload

with open(TKN) as f: cr = Credentials.from_authorized_user_info(json.load(f), scopes=["https://www.googleapis.com/auth/drive"])
svc = build("drive", "v3", credentials=cr)
for nm in ["Custos_Projecao_OmniChannel.pptx"]:
    old = svc.files().list(q=f"name='{nm}' and '{FID}' in parents and trashed=false", fields="files(id)").execute()
    for f in old.get("files", []): svc.files().delete(fileId=f["id"]).execute()
# Upload dated version to Custos/
old = svc.files().list(q=f"name='{FNAME}' and '{FID}' in parents and trashed=false", fields="files(id)").execute()
for f in old.get("files", []): svc.files().delete(fileId=f["id"]).execute()
up = svc.files().create(body={"name": FNAME, "parents": [FID]},
    media_body=MediaFileUpload(OUT, mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    fields="id,webViewLink").execute()
print(f"  Drive: {up['webViewLink']}")
