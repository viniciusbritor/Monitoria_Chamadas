"""
Add "Hoje" slide + PIX QR Code to existing PPTX
Uses python-pptx for local edits (Layer 3)
"""
import json, os, io
import qrcode
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── 1. Generate PIX QR Code ───
CPF = "04779977754"
NAME = "VINICIUS BRITO"
CITY = "SAO PAULO"
PIX_KEY_TYPE = "CPF"

def generate_pix_payload(key, name, city, amount="0"):
    """Generate PIX EMV BR Code payload string."""
    # CRC16 function
    def crc16_ccitt(data):
        """CRC16-CCITT (0x1021) with initial value 0xFFFF."""
        crc = 0xFFFF
        for byte in data:
            crc ^= (byte << 8)
            for _ in range(8):
                if crc & 0x8000:
                    crc = ((crc << 1) ^ 0x1021) & 0xFFFF
                else:
                    crc = (crc << 1) & 0xFFFF
        return crc
    
    payload = ""
    # 00 - Payload Format Indicator
    payload += "000201"
    # 01 - Point of Initiation (static)
    # skipped for simplicity
    # 26 - Merchant Account Information
    gui = "br.gov.bcb.pix"
    gui_block = f"0014{gui}"
    
    # CPF key is just the number
    key_block = f"01{len(key):02d}{key}"
    merchant_info = f"{gui_block}{key_block}"
    payload += f"26{len(merchant_info):02d}{merchant_info}"
    
    # 52 - MCC
    payload += "52040000"
    # 53 - Currency (986 = BRL)
    payload += "5303986"
    # 54 - Amount
    if amount and amount != "0":
        payload += f"54{len(amount):02d}{amount}"
    # 58 - Country
    payload += "5802BR"
    # 59 - Merchant Name
    payload += f"59{len(name):02d}{name}"
    # 60 - City
    payload += f"60{len(city):02d}{city}"
    # 62 - Additional Data (reference)
    ref = "***"
    ref_block = f"05{len(ref):02d}{ref}"
    payload += f"62{len(ref_block):02d}{ref_block}"
    # 63 - CRC16
    payload += "6304"
    crc = crc16_ccitt(payload.encode('ascii'))
    payload += f"{crc:04X}"
    
    return payload

pix_payload = generate_pix_payload(CPF, NAME, CITY)
print(f"PIX payload ({len(pix_payload)} chars): {pix_payload[:80]}...")

# Generate QR code image
qr = qrcode.QRCode(version=2, box_size=10, border=4)
qr.add_data(pix_payload)
qr.make(fit=True)
qr_img = qr.make_image(fill_color="black", back_color="white")

QR_PATH = r"C:\Users\vinic\workspace_antigravity\Monitoria_Chamadas\docs\pix_qr.png"
qr_img.save(QR_PATH)
print(f"QR code saved: {QR_PATH}")

# ─── 2. Add slides to existing PPTX ───
PPTX_PATH = r"C:\Users\vinic\workspace_antigravity\Monitoria_Chamadas\docs\Custos_Projecao.pptx"

prs = Presentation(PPTX_PATH)
layout = prs.slide_layouts[0]  # first available layout
print(f"Layouts available: {len(prs.slide_layouts)}")

# ─── SLIDE: Custos Reais de Hoje ───
slide_today = prs.slides.add_slide(layout)

# Title
left, top, width, height = Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.8)
txBox = slide_today.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "CENARIO ATUAL (HOJE) — Infra em Operacao"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
p.alignment = PP_ALIGN.CENTER

# Today's costs table
table_data = [
    ["Componente", "Custo/mes (USD)", "Custo/mes (BRL)", "Nota"],
    ["Worker (min=1, sempre ligado)", "$77", "R$ 424", "4vCPU/4GB, 730h/mes fixo"],
    ["DeepSeek LLM (poucas chamadas)", "$2", "R$ 11", "10-50 chamadas/dia reais"],
    ["API Cloud Run (idle)", "$1", "R$ 6", "Min=0, quase nao ativa"],
    ["Firestore + Storage + PubSub", "$15", "R$ 83", "Infra fixa compartilhada"],
    ["TOTAL MONITORIA", "$95", "R$ 524", " "],
    ["5 Chatbots WhatsApp", "$113", "R$ 622", "Cloud Run + Postgres + LLM"],
    ["TOTAL ECOSSISTEMA HOJE", "$208", "R$ 1.146", "Monitoria + WhatsApp"],
]

rows, cols = len(table_data), len(table_data[0])
tbl_shape = slide_today.shapes.add_table(rows, cols, Inches(0.5), Inches(1.6), Inches(9), Inches(0.4 * rows))
tbl = tbl_shape.table

for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = table_data[r][c]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            if r == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            elif r in [5, 7]:
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(0x00, 0x61, 0x00)
            else:
                paragraph.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            paragraph.alignment = PP_ALIGN.CENTER
        
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        elif r in [5, 7]:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE2, 0xEF, 0xDA)

# Capacity note
left2, top2 = Inches(1.5), Inches(5.8)
txBox2 = slide_today.shapes.add_textbox(left2, top2, Inches(7), Inches(1.2))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "Capacidade HOJE: ~500-1.000 chamadas/dia (sem gargalo)"
p2.font.size = Pt(14)
p2.font.bold = True
p2.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
p2.alignment = PP_ALIGN.CENTER

p3 = tf2.add_paragraph()
p3.text = "1 instancia (min=1) = ~170 chamadas/hora | Pico (max=4) = ~680 chamadas/hora"
p3.font.size = Pt(10)
p3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
p3.alignment = PP_ALIGN.CENTER

# ─── SLIDE: PIX QR Code ───
slide_pix = prs.slides.add_slide(layout)

# Title
txBox3 = slide_pix.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
tf3 = txBox3.text_frame
p4 = tf3.paragraphs[0]
p4.text = "Apoie o Projeto OmniChannel"
p4.font.size = Pt(32)
p4.font.bold = True
p4.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
p4.alignment = PP_ALIGN.CENTER

# Subtitle
p5 = tf3.add_paragraph()
p5.text = "Contribua via PIX para o desenvolvimento contínuo"
p5.font.size = Pt(14)
p5.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
p5.alignment = PP_ALIGN.CENTER

# QR Code image
qr_left = Inches(3.2)
qr_top = Inches(2.0)
qr_height = Inches(3.5)
slide_pix.shapes.add_picture(QR_PATH, qr_left, qr_top, height=qr_height)

# CPF info
txBox4 = slide_pix.shapes.add_textbox(Inches(1.5), Inches(5.8), Inches(7), Inches(0.8))
tf4 = txBox4.text_frame
p6 = tf4.paragraphs[0]
p6.text = "PIX CPF: 047.799.777-54"
p6.font.size = Pt(20)
p6.font.bold = True
p6.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
p6.alignment = PP_ALIGN.CENTER

p7 = tf4.add_paragraph()
p7.text = "Vinícius Brito | OmniChannel — Projeções Financeiras 2026"
p7.font.size = Pt(10)
p7.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
p7.alignment = PP_ALIGN.CENTER

# ─── SAVE ───
OUTPUT = r"C:\Users\vinic\workspace_antigravity\Monitoria_Chamadas\docs\Custos_Projecao_Final.pptx"
prs.save(OUTPUT)
print(f"Updated PPTX saved: {OUTPUT}")

# ─── Upload to Drive ───
import sys
sys.path.insert(0, os.path.dirname(__file__))
from google.oauth2.credentials import Credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload

TOKEN = os.path.expanduser(r"~\.gemini\config\skills\google_calendar_manager\resources\token_drive.json")
FOLDER = "1aNCHHOiQQzquuxzaeQQa8qr3ciZcsfMt"

with open(TOKEN) as f:
    creds = Credentials.from_authorized_user_info(json.load(f), scopes=["https://www.googleapis.com/auth/drive"])
service = build("drive", "v3", credentials=creds)

# Delete old
for name in ["Custos_Projecao_Final.pptx", "Custos_Projecao.pptx"]:
    old = service.files().list(q=f"name='{name}' and '{FOLDER}' in parents and trashed=false", fields="files(id)").execute()
    for f in old.get("files", []):
        service.files().delete(fileId=f["id"]).execute()

meta = {"name": "Custos_Projecao_Final.pptx", "parents": [FOLDER]}
media = MediaFileUpload(OUTPUT, mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")
up = service.files().create(body=meta, media_body=media, fields="id,webViewLink").execute()
print(f"Drive: {up['webViewLink']}")
