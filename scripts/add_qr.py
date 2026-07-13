"""Add PIX QR code to slide 10 of the v2 PPTX"""
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from google.oauth2.credentials import Credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload

PPTX_IN = r"C:\Users\vinic\workspace_antigravity\Monitoria_Chamadas\docs\Custos_Projecao_v2.pptx"
PPTX_OUT = r"C:\Users\vinic\workspace_antigravity\Monitoria_Chamadas\docs\Custos_Projecao_OmniChannel.pptx"
QR_PATH = r"C:\Users\vinic\workspace_antigravity\Monitoria_Chamadas\docs\pix_qr.png"

prs = Presentation(PPTX_IN)

# Add new slide with PIX QR Code
layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(layout)

# QR Code image
slide.shapes.add_picture(QR_PATH, Inches(3.0), Inches(1.2), height=Inches(3.8))

# CPF info below QR
txBox = slide.shapes.add_textbox(Inches(2), Inches(5.3), Inches(6), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "PIX CPF: 047.799.777-54"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Vinícius Brito · OmniChannel 2026"
p2.font.size = Pt(11)
p2.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
p2.alignment = PP_ALIGN.CENTER

prs.save(PPTX_OUT)
print(f"Saved: {PPTX_OUT}")

# Upload to Drive
TOKEN = os.path.expanduser(r"~\.gemini\config\skills\google_calendar_manager\resources\token_drive.json")
FOLDER = "1aNCHHOiQQzquuxzaeQQa8qr3ciZcsfMt"

with open(TOKEN) as f:
    creds = Credentials.from_authorized_user_info(json.load(f), scopes=["https://www.googleapis.com/auth/drive"])
service = build("drive", "v3", credentials=creds)

for name in ["Custos_Projecao_OmniChannel.pptx"]:
    old = service.files().list(q=f"name='{name}' and '{FOLDER}' in parents and trashed=false", fields="files(id)").execute()
    for f in old.get("files", []):
        service.files().delete(fileId=f["id"]).execute()

meta = {"name": "Custos_Projecao_OmniChannel.pptx", "parents": [FOLDER]}
media = MediaFileUpload(PPTX_OUT, mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")
up = service.files().create(body=meta, media_body=media, fields="id,webViewLink").execute()
print(f"Drive: {up['webViewLink']}")
print(f"Local: {PPTX_OUT}")
